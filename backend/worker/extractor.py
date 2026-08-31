"""Deterministic, local extraction for presentation material.

This module deliberately has no database or network dependency.  It is safe to
use in the worker process and gives a useful semantic result when converters,
OCR, or an AI provider are unavailable.
"""
from __future__ import annotations

import hashlib
import io
import json
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import PurePosixPath
from typing import Any

MAX_BYTES = 25 * 1024 * 1024
MAX_ZIP_MEMBERS = 2_000
MAX_ZIP_UNCOMPRESSED = 150 * 1024 * 1024
MAX_PAGES = 100
MAX_SLIDES = 60


class MaterialError(ValueError):
    pass


@dataclass
class ExtractedMaterial:
    units: list[dict[str, Any]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


def _stable_id(unit_key: str, kind: str, ordinal: int) -> str:
    return hashlib.sha256(f"{unit_key}:{kind}:{ordinal}".encode()).hexdigest()[:20]


def _text(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def validate_source(data: bytes, source_type: str) -> None:
    if len(data) > MAX_BYTES:
        raise MaterialError("source exceeds the 25 MB limit")
    if source_type in {"pptx", "docx"}:
        _validate_zip(data)
    if source_type in {"ppt", "doc"}:
        # OLE encrypted files cannot be inspected safely here; LibreOffice will
        # reject encrypted/malformed payloads during conversion.
        if data[:8] != bytes.fromhex("D0CF11E0A1B11AE1"):
            raise MaterialError("malformed legacy Office file")


def _validate_zip(data: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            members = archive.infolist()
            if len(members) > MAX_ZIP_MEMBERS:
                raise MaterialError("archive has too many members")
            total = 0
            for member in members:
                path = PurePosixPath(member.filename)
                if member.filename.startswith(("/", "\\")) or ".." in path.parts:
                    raise MaterialError("archive contains an unsafe path")
                if member.file_size > MAX_BYTES:
                    raise MaterialError("Office archive contains an oversized embedded asset")
                total += member.file_size
                if total > MAX_ZIP_UNCOMPRESSED:
                    raise MaterialError("archive expands beyond safe limit")
                if member.flag_bits & 0x1:
                    raise MaterialError("encrypted Office files are unsupported")
    except zipfile.BadZipFile as exc:
        raise MaterialError("malformed Office archive") from exc


def _analysis(text: str) -> dict[str, Any]:
    numbers = sorted(set(re.findall(r"(?<![\w.])-?\d+(?:\.\d+)?%?", text)))[:20]
    words = re.findall(r"[A-Za-z][A-Za-z0-9-]{3,}", text.lower())
    ignored = {"this", "that", "with", "from", "have", "will", "your", "about", "their", "there", "which"}
    frequency: dict[str, int] = {}
    for word in words:
        if word not in ignored:
            frequency[word] = frequency.get(word, 0) + 1
    concepts = [word for word, _ in sorted(frequency.items(), key=lambda item: (-item[1], item[0]))[:3]]
    return {"facts": numbers, "concepts": concepts}


def _unit(unit_key: str, ordinal: int, unit_type: str, title: str, elements: list[dict[str, Any]], notes: str = "") -> dict[str, Any]:
    for element in elements:
        if isinstance(element, dict):
            element.setdefault("source_class", str(element.get("provenance") or "native"))
    text = " ".join(_text(str(element.get("text", ""))) for element in elements if element.get("text"))
    return {
        "unit_key": unit_key,
        # The database contract deliberately reserves zero: ordinal is the
        # human-facing page/slide number and has CHECK (ordinal > 0).
        "ordinal": ordinal + 1,
        "unit_type": unit_type,
        "title": _text(title) or f"{unit_type.title()} {ordinal + 1}",
        "content": {"elements": elements, "provenance": "native"},
        "notes": _text(notes) or None,
        "analysis": _analysis(text),
        "search_text": _text(" ".join((title, text, notes))),
    }


def extract_pptx(data: bytes) -> ExtractedMaterial:
    from pptx import Presentation

    try:
        deck = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise MaterialError("malformed or encrypted PPTX") from exc
    if len(deck.slides) > MAX_SLIDES:
        raise MaterialError("presentation exceeds the 60 slide limit")
    result = ExtractedMaterial(warnings=_pptx_visual_warnings(data))
    notes = _pptx_notes(data)
    for ordinal, slide in enumerate(deck.slides):
        key = f"slide-{ordinal + 1}"
        elements: list[dict[str, Any]] = []
        title = ""
        for position, shape in enumerate(slide.shapes):
            bounds = {"x": int(shape.left), "y": int(shape.top), "width": int(shape.width), "height": int(shape.height)}
            element: dict[str, Any] = {"id": _stable_id(key, "shape", position), "provenance": "native", "bounds": bounds}
            if getattr(shape, "has_text_frame", False):
                paragraphs = [(paragraph_index, paragraph, _text(paragraph.text)) for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs) if _text(paragraph.text)]
                if paragraphs:
                    if not title and (getattr(shape, "is_placeholder", False) or position == 0):
                        title = paragraphs[0][2]
                    for paragraph_index, paragraph, value in paragraphs:
                        # List hierarchy is a semantic signal. Keeping each
                        # paragraph separately preserves ordering and nesting.
                        elements.append({**element, "id": _stable_id(key, "text", position * 10_000 + paragraph_index), "type": "text", "text": value, "level": int(paragraph.level), "order": paragraph_index})
            if getattr(shape, "has_table", False):
                rows = [[_text(cell.text) for cell in row.cells] for row in shape.table.rows]
                elements.append({**element, "id": _stable_id(key, "table", position), "type": "table", "rows": rows, "text": " ".join(" ".join(row) for row in rows)})
            if getattr(shape, "shape_type", None) is not None and str(getattr(shape, "shape_type", "")).lower().find("picture") >= 0:
                elements.append({**element, "id": _stable_id(key, "picture", position), "type": "picture", "alt_text": getattr(shape, "name", "picture")})
            if getattr(shape, "has_chart", False):
                elements.append({**element, "id": _stable_id(key, "chart", position), "type": "chart", "chart_type": str(shape.chart.chart_type)})
        result.units.append(_unit(key, ordinal, "slide", title, elements, notes.get(ordinal, "")))
    return result


def _pptx_visual_warnings(data: bytes) -> list[str]:
    """Flag OOXML constructs whose semantics are not faithfully exposed by python-pptx."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = archive.namelist()
            warnings = []
            if any(name.startswith("ppt/diagrams/") for name in names):
                warnings.append("SmartArt detected; native semantic extraction may be incomplete")
            if any(name.startswith("ppt/embeddings/") for name in names):
                warnings.append("Embedded object detected; its content was not semantically extracted")
            media_extensions = {".mp4", ".mov", ".avi", ".mp3", ".wav", ".m4a"}
            if any(
                name.startswith("ppt/media/") and PurePosixPath(name).suffix.lower() in media_extensions
                for name in names
            ):
                warnings.append("Embedded audio or video detected; playback was not extracted")
            if any(b"<p:timing" in archive.read(name) for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)):
                warnings.append("Slide animations detected; animation sequence was not extracted")
            return warnings
    except Exception:
        return []


def _pptx_notes(data: bytes) -> dict[int, str]:
    """Read note-slide text directly because python-pptx exposes it unevenly."""
    import xml.etree.ElementTree as ET
    found: dict[int, str] = {}
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            for name in archive.namelist():
                match = re.fullmatch(r"ppt/notesSlides/notesSlide(\d+)\.xml", name)
                if match:
                    root = ET.fromstring(archive.read(name))
                    found[int(match.group(1)) - 1] = " ".join(node.text or "" for node in root.iter() if node.tag.endswith("}t"))
    except Exception:
        return {}
    return found


def extract_docx(data: bytes) -> ExtractedMaterial:
    from docx import Document

    try:
        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise MaterialError("malformed or encrypted DOCX") from exc
    units: list[dict[str, Any]] = []
    current_title = "Document"
    elements: list[dict[str, Any]] = []
    ordinal = 0

    def flush() -> None:
        nonlocal ordinal, elements
        if elements:
            units.append(_unit(f"section-{ordinal + 1}", ordinal, "section", current_title, elements))
            ordinal += 1
            elements = []

    for paragraph in document.paragraphs:
        value = _text(paragraph.text)
        if not value:
            continue
        if paragraph.style and paragraph.style.name.lower().startswith("heading"):
            flush()
            current_title = value
        else:
            elements.append({"id": _stable_id(f"section-{ordinal + 1}", "paragraph", len(elements)), "type": "paragraph", "text": value, "provenance": "native"})
    for table in document.tables:
        rows = [[_text(cell.text) for cell in row.cells] for row in table.rows]
        elements.append({"id": _stable_id(f"section-{ordinal + 1}", "table", len(elements)), "type": "table", "rows": rows, "text": " ".join(" ".join(row) for row in rows), "provenance": "native"})
    flush()
    if len(units) > MAX_PAGES:
        raise MaterialError("document exceeds the 100 section limit")
    return ExtractedMaterial(units=units or [_unit("section-1", 0, "section", "Document", [])])


def extract_txt(data: bytes) -> ExtractedMaterial:
    text = data.decode("utf-8-sig", errors="replace")
    chunks = [chunk.strip() for chunk in re.split(r"\n\s*\n+", text) if chunk.strip()]
    if len(chunks) > MAX_PAGES:
        raise MaterialError("text exceeds the 100 section limit")
    return ExtractedMaterial(units=[_unit(f"section-{i + 1}", i, "section", chunk.splitlines()[0][:120], [{"id": _stable_id(f"section-{i + 1}", "paragraph", 0), "type": "paragraph", "text": chunk, "provenance": "native"}]) for i, chunk in enumerate(chunks)] or [_unit("section-1", 0, "section", "Text", [])])


def extract_pdf(data: bytes) -> ExtractedMaterial:
    import pdfplumber

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            if len(pdf.pages) > MAX_PAGES:
                raise MaterialError("PDF exceeds the 100 page limit")
            units = []
            for ordinal, page in enumerate(pdf.pages):
                key = f"page-{ordinal + 1}"
                text = _text(page.extract_text() or "")
                elements = []
                if text:
                    elements.append({
                        "id": _stable_id(key, "text", 0),
                        "type": "text",
                        "text": text,
                        "provenance": "native",
                        "bounds": {"x": 0, "y": 0, "width": float(page.width), "height": float(page.height)},
                    })
                try:
                    for table_index, table in enumerate(page.find_tables()[:10]):
                        rows = [
                            [_text(str(cell or "")) for cell in row]
                            for row in (table.extract() or [])[:100]
                        ]
                        if not any(any(cell for cell in row) for row in rows):
                            continue
                        x0, top, x1, bottom = table.bbox
                        elements.append({
                            "id": _stable_id(key, "table", table_index),
                            "type": "table",
                            "rows": rows,
                            "text": " ".join(" ".join(row) for row in rows),
                            "provenance": "native",
                            "bounds": {"x": x0, "y": top, "width": x1 - x0, "height": bottom - top},
                        })
                except Exception:
                    # Page text remains usable when table geometry is malformed.
                    pass
                units.append(_unit(key, ordinal, "page", f"Page {ordinal + 1}", elements))
            return ExtractedMaterial(units=units)
    except MaterialError:
        raise
    except Exception as exc:
        raise MaterialError("malformed or encrypted PDF") from exc


def extract(data: bytes, source_type: str) -> ExtractedMaterial:
    validate_source(data, source_type)
    handlers = {"pptx": extract_pptx, "docx": extract_docx, "txt": extract_txt, "pdf": extract_pdf}
    if source_type not in handlers:
        raise MaterialError(f"{source_type.upper()} requires LibreOffice conversion")
    return handlers[source_type](data)


def summarize(units: list[dict[str, Any]]) -> dict[str, Any]:
    facts = sorted({fact for unit in units for fact in unit["analysis"].get("facts", [])})
    concepts = sorted({concept for unit in units for concept in unit["analysis"].get("concepts", [])})[:20]
    return {"fact_ledger": facts, "concepts": concepts, "unit_count": len(units), "generator": "deterministic"}
